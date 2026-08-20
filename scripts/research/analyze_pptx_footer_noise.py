#!/usr/bin/env python3
"""对本地 PPTX 中普通文本框的跨页重复模式做脱敏聚合。

脚本不会输出文件名或原文。它把每份输入匿名编号，并只输出数量、位置类别、
长度区间和 AnyDoc 是否保留候选文本等聚合信息。
"""

from __future__ import annotations

import argparse
import io
import json
import math
import posixpath
import re
import unicodedata
import zipfile
from collections import Counter, defaultdict
from dataclasses import dataclass
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path

import anydoc
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", unicodedata.normalize("NFKC", text)).strip()


def text_units(shape) -> list[str]:
    units = [normalize(paragraph.text) for paragraph in shape.text_frame.paragraphs]
    units = [unit for unit in units if unit]
    return units or ([normalize(shape.text)] if normalize(shape.text) else [])


def band(top: float, height: float) -> str:
    center = top + height / 2
    if center >= 0.82:
        return "bottom"
    if center <= 0.18:
        return "top"
    return "middle"


def length_bucket(text: str) -> str:
    size = len(text)
    if size <= 12:
        return "1-12"
    if size <= 40:
        return "13-40"
    if size <= 100:
        return "41-100"
    return "101+"


def iter_inline_text(inlines) -> list[str]:
    values: list[str] = []
    for inline in inlines or []:
        if inline.kind == "text" and inline.text:
            values.append(inline.text)
        elif inline.kind == "link":
            values.extend(iter_inline_text(inline.content))
        elif inline.kind == "image" and inline.alt:
            values.append(inline.alt)
        elif inline.kind == "line_break":
            values.append("\n")
    return values


def iter_block_text(blocks) -> list[str]:
    values: list[str] = []
    for block in blocks or []:
        if block.kind in {"heading", "paragraph"}:
            values.extend(iter_inline_text(block.content))
        elif block.kind == "list" and block.list:
            for item in block.list.items:
                values.extend(iter_block_text(item.blocks))
        elif block.kind == "table" and block.table:
            for row in block.table.grid:
                for slot in row:
                    if slot.kind == "origin" and slot.cell:
                        values.extend(iter_block_text(slot.cell.blocks))
        elif block.kind == "block_quote":
            values.extend(iter_block_text(block.blocks))
        elif block.kind == "code_block" and block.text:
            values.append(block.text)
    return values


def relationship_part(part: str) -> str:
    directory, filename = posixpath.split(part)
    return posixpath.join(directory, "_rels", f"{filename}.rels")


def relationship_targets(src: zipfile.ZipFile, part: str) -> list[str]:
    rels_part = relationship_part(part)
    if rels_part not in src.namelist():
        return []
    rels = etree.fromstring(src.read(rels_part))
    targets: list[str] = []
    for rel in rels:
        if rel.get("TargetMode") == "External":
            continue
        rel_type = rel.get("Type", "")
        # slideMaster 会反向列出所有 layout；所选 layout 已由 slide 的关系进入闭包。
        if "/slideMasters/" in f"/{part}" and rel_type.endswith("/slideLayout"):
            continue
        target = rel.get("Target")
        if target:
            targets.append(posixpath.normpath(posixpath.join(posixpath.dirname(part), target)))
    return targets


def project_single_slide(source: Path, slide_index: int) -> bytes:
    """收窄 sldIdLst，并保留所选页关系闭包与演示级公共部件。"""
    output = io.BytesIO()
    with zipfile.ZipFile(source) as src:
        presentation = etree.fromstring(src.read("ppt/presentation.xml"))
        slide_ids = presentation.xpath("//p:sldIdLst/p:sldId", namespaces={"p": P_NS})
        if slide_index >= len(slide_ids):
            raise IndexError(slide_index)
        selected = slide_ids[slide_index]
        selected_rid = selected.get(f"{{{R_NS}}}id")
        parent = selected.getparent()
        for index, slide_id in enumerate(list(parent)):
            if index != slide_index:
                parent.remove(slide_id)
        projected_xml = etree.tostring(
            presentation, xml_declaration=True, encoding="UTF-8", standalone=True
        )
        presentation_rels = etree.fromstring(src.read("ppt/_rels/presentation.xml.rels"))
        selected_part = None
        for rel in presentation_rels:
            target = rel.get("Target")
            if not target or rel.get("TargetMode") == "External":
                continue
            rel_type = rel.get("Type", "")
            if rel_type.endswith("/slide") and rel.get("Id") == selected_rid:
                selected_part = posixpath.normpath(posixpath.join("ppt", target))
                break
        if not selected_part:
            raise ValueError("selected slide relationship is missing")

        keep = {
            "[Content_Types].xml",
            "_rels/.rels",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        }
        visited: set[str] = set()
        pending = [selected_part]
        while pending:
            part = pending.pop()
            if part in visited or part not in src.namelist():
                continue
            visited.add(part)
            keep.add(part)
            rels_part = relationship_part(part)
            if rels_part in src.namelist():
                keep.add(rels_part)
                pending.extend(relationship_targets(src, part))

        with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as dst:
            for name in sorted(keep):
                payload = projected_xml if name == "ppt/presentation.xml" else src.read(name)
                dst.writestr(name, payload)
    return output.getvalue()


@dataclass(frozen=True)
class Occurrence:
    page: int
    top: float
    left: float
    width: float
    height: float
    position_band: str
    rotation: float
    anydoc_retained: bool


def iter_text_shapes(shapes, slide_width: int, slide_height: int, transform=None):
    """递归展开组合图形，并把子坐标映射回页坐标。"""
    for shape in shapes:
        if transform is None:
            abs_left, abs_top = float(shape.left), float(shape.top)
            abs_width, abs_height = float(shape.width), float(shape.height)
        else:
            outer_left, outer_top, outer_width, outer_height, ch_left, ch_top, ch_width, ch_height = transform
            abs_left = outer_left + (float(shape.left) - ch_left) * outer_width / ch_width
            abs_top = outer_top + (float(shape.top) - ch_top) * outer_height / ch_height
            abs_width = float(shape.width) * outer_width / ch_width
            abs_height = float(shape.height) * outer_height / ch_height

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            xfrm = shape._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
            ch_off = xfrm.find(qn("a:chOff"))
            ch_ext = xfrm.find(qn("a:chExt"))
            child_transform = (
                abs_left,
                abs_top,
                abs_width,
                abs_height,
                float(ch_off.get("x")),
                float(ch_off.get("y")),
                float(ch_ext.get("cx")),
                float(ch_ext.get("cy")),
            )
            yield from iter_text_shapes(shape.shapes, slide_width, slide_height, child_transform)
        elif getattr(shape, "has_text_frame", False):
            yield (
                shape,
                abs_left / slide_width,
                abs_top / slide_height,
                abs_width / slide_width,
                abs_height / slide_height,
            )


def is_position_stable(items: list[Occurrence], tolerance: float = 0.02) -> bool:
    return (
        max(x.top for x in items) - min(x.top for x in items) <= tolerance
        and max(x.left for x in items) - min(x.left for x in items) <= tolerance
        and max(x.width for x in items) - min(x.width for x in items) <= tolerance
        and max(x.height for x in items) - min(x.height for x in items) <= tolerance
    )


def analyze_document(source: Path, anonymous_id: str) -> dict:
    prs = Presentation(source)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    page_count = len(prs.slides)
    occurrences: dict[str, list[Occurrence]] = defaultdict(list)
    all_page_signatures: list[tuple[str, ...]] = []
    anydoc_block_kinds: Counter[str] = Counter()
    conversion_failures = 0
    standard_noise_placeholders = 0
    standard_noise_placeholders_retained = 0

    for page_index, slide in enumerate(prs.slides):
        projected = project_single_slide(source, page_index)
        try:
            document = anydoc.to_document(projected, format="pptx")
            anydoc_block_kinds.update(block.kind for block in document.blocks)
            anydoc_text = normalize(" ".join(iter_block_text(document.blocks)))
        except Exception:
            conversion_failures += 1
            anydoc_text = ""

        page_texts: list[str] = []
        for shape, left, top, width, height in iter_text_shapes(
            slide.shapes, slide_width, slide_height
        ):
            units = text_units(shape)
            if not units:
                continue
            page_texts.extend(units)
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type in {
                    PP_PLACEHOLDER.DATE,
                    PP_PLACEHOLDER.FOOTER,
                    PP_PLACEHOLDER.SLIDE_NUMBER,
                }:
                    standard_noise_placeholders += len(units)
                    standard_noise_placeholders_retained += sum(
                        bool(anydoc_text and unit in anydoc_text) for unit in units
                    )
                continue
            for text in units:
                occurrences[text].append(
                    Occurrence(
                        page=page_index,
                        top=top,
                        left=left,
                        width=width,
                        height=height,
                        position_band=band(top, height),
                        rotation=float(shape.rotation or 0),
                        anydoc_retained=bool(anydoc_text and text in anydoc_text),
                    )
                )
        all_page_signatures.append(tuple(sorted(set(page_texts))))

    signature_counts = Counter(all_page_signatures)
    duplicate_signature_groups = sum(count >= 2 for count in signature_counts.values())
    duplicate_signature_pages = sum(count for count in signature_counts.values() if count >= 2)
    repeated_groups: list[dict] = []
    for text, raw_items in occurrences.items():
        # 同一页重复出现多次时，只按页计覆盖率，保留首次几何位置。
        items_by_page: dict[int, Occurrence] = {}
        for item in raw_items:
            items_by_page.setdefault(item.page, item)
        items = list(items_by_page.values())
        if len(items) < 2:
            continue
        bands = Counter(item.position_band for item in items)
        dominant_band, dominant_band_count = bands.most_common(1)[0]
        pages = {item.page for item in items}
        duplicate_page_count = sum(signature_counts[all_page_signatures[p]] > 1 for p in pages)
        repeated_groups.append(
            {
                "page_occurrences": len(items),
                "prevalence": round(len(items) / page_count, 4),
                "position_band": dominant_band,
                "band_consistency": round(dominant_band_count / len(items), 4),
                "position_stable": is_position_stable(items),
                "length_bucket": length_bucket(text),
                "multiline": "\n" in text,
                "rotated": any(abs(item.rotation) >= 5 for item in items),
                "all_retained_by_anydoc": all(item.anydoc_retained for item in items),
                "retained_occurrences": sum(item.anydoc_retained for item in items),
                "duplicate_slide_share": round(duplicate_page_count / len(items), 4),
            }
        )

    strict_candidates = [
        group
        for group in repeated_groups
        if group["page_occurrences"] >= max(5, math.ceil(page_count * 0.5))
        and group["position_band"] == "bottom"
        and group["band_consistency"] == 1
        and group["position_stable"]
        and group["length_bucket"] in {"1-12", "13-40"}
    ]
    broad_candidates = [
        group
        for group in repeated_groups
        if group["page_occurrences"] >= max(3, math.ceil(page_count * 0.2))
        and group["position_stable"]
    ]
    return {
        "document": anonymous_id,
        "pages": page_count,
        "conversion_failures": conversion_failures,
        "standard_noise_placeholder_text_units": standard_noise_placeholders,
        "standard_noise_placeholder_units_retained": standard_noise_placeholders_retained,
        "anydoc_top_level_block_kinds": dict(sorted(anydoc_block_kinds.items())),
        "ordinary_text_distinct": len(occurrences),
        "repeated_exact_groups_2plus": len(repeated_groups),
        "repeated_exact_groups_3plus": sum(x["page_occurrences"] >= 3 for x in repeated_groups),
        "repeated_exact_group_occurrences": sum(x["page_occurrences"] for x in repeated_groups),
        "position_stable_repeated_groups": sum(x["position_stable"] for x in repeated_groups),
        "duplicate_text_signature_groups": duplicate_signature_groups,
        "duplicate_text_signature_pages": duplicate_signature_pages,
        "repeated_group_bands": dict(sorted(Counter(x["position_band"] for x in repeated_groups).items())),
        "repeated_group_length_buckets": dict(
            sorted(Counter(x["length_bucket"] for x in repeated_groups).items())
        ),
        "repeated_groups_all_retained_by_anydoc": sum(
            x["all_retained_by_anydoc"] for x in repeated_groups
        ),
        "broad_stable_candidates": len(broad_candidates),
        "broad_candidate_bands": dict(
            sorted(Counter(x["position_band"] for x in broad_candidates).items())
        ),
        "strict_footer_candidates": len(strict_candidates),
        "strict_footer_candidate_occurrences": sum(
            x["page_occurrences"] for x in strict_candidates
        ),
        "strict_footer_candidates_all_retained": sum(
            x["all_retained_by_anydoc"] for x in strict_candidates
        ),
        "strict_footer_examples": [
            {
                "page_occurrences": x["page_occurrences"],
                "of_pages": page_count,
                "position": x["position_band"],
                "length_bucket": x["length_bucket"],
                "retained_by_anydoc": x["all_retained_by_anydoc"],
                "duplicate_slide_share": x["duplicate_slide_share"],
            }
            for x in strict_candidates[:5]
        ],
        "broad_candidate_examples": [
            {
                "page_occurrences": x["page_occurrences"],
                "of_pages": page_count,
                "position": x["position_band"],
                "length_bucket": x["length_bucket"],
                "retained_by_anydoc": x["all_retained_by_anydoc"],
                "duplicate_slide_share": x["duplicate_slide_share"],
            }
            for x in broad_candidates[:10]
        ],
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("fixtures", type=Path)
    args = parser.parse_args()
    sources = sorted(args.fixtures.glob("*.pptx"))
    documents = [
        analyze_document(path, f"doc-{index:02d}")
        for index, path in enumerate(sources, start=1)
    ]
    try:
        anydoc_version = version("firecrawl-anydoc")
    except PackageNotFoundError:
        anydoc_version = "unknown"
    summary = {
        "analyzer": {
            "anydoc_distribution": "firecrawl-anydoc",
            "anydoc_version": anydoc_version,
            "call": "to_document(projected_bytes, format='pptx')",
            "projection": "sldIdLst narrowed to one slide; selected slide relationship closure retained",
        },
        "totals": {
            "documents": len(documents),
            "pages": sum(x["pages"] for x in documents),
            "conversion_failures": sum(x["conversion_failures"] for x in documents),
            "standard_noise_placeholder_text_units": sum(
                x["standard_noise_placeholder_text_units"] for x in documents
            ),
            "standard_noise_placeholder_units_retained": sum(
                x["standard_noise_placeholder_units_retained"] for x in documents
            ),
            "repeated_exact_groups_2plus": sum(x["repeated_exact_groups_2plus"] for x in documents),
            "repeated_exact_groups_3plus": sum(x["repeated_exact_groups_3plus"] for x in documents),
            "repeated_exact_group_occurrences": sum(
                x["repeated_exact_group_occurrences"] for x in documents
            ),
            "position_stable_repeated_groups": sum(
                x["position_stable_repeated_groups"] for x in documents
            ),
            "duplicate_text_signature_groups": sum(
                x["duplicate_text_signature_groups"] for x in documents
            ),
            "duplicate_text_signature_pages": sum(
                x["duplicate_text_signature_pages"] for x in documents
            ),
            "broad_stable_candidates": sum(x["broad_stable_candidates"] for x in documents),
            "strict_footer_candidates": sum(x["strict_footer_candidates"] for x in documents),
            "strict_footer_candidate_occurrences": sum(
                x["strict_footer_candidate_occurrences"] for x in documents
            ),
            "strict_footer_candidates_all_retained": sum(
                x["strict_footer_candidates_all_retained"] for x in documents
            ),
        },
        "documents": documents,
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
