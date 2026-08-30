from __future__ import annotations

from collections.abc import Iterator
from dataclasses import replace
from io import BytesIO
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from pptextract.api import create_app
from pptextract.config import Settings
from pptextract.conversion import NormalizedImage, NormalizedPageContent
from pptextract.db import connect, transaction
from pptextract.pptx_projection import SourcePage
from pptextract.rendering import StandardPageRender
from pptextract.worker import run_once

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture
def system(tmp_path: Path) -> Iterator[tuple[TestClient, Settings]]:
    settings = replace(Settings.for_test(tmp_path), job_retry_base_seconds=0)
    with TestClient(create_app(settings)) as client:
        yield client, settings


def _presentation(*titles: str, subject: str = "") -> bytes:
    presentation = Presentation()
    presentation.core_properties.subject = subject
    for title in titles:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _remove_title(source: bytes, title: str) -> bytes:
    presentation = Presentation(BytesIO(source))
    index = next(
        index
        for index, slide in enumerate(presentation.slides)
        if slide.shapes.title.text == title
    )
    slide_id = presentation.slides._sldIdLst[index]
    presentation.part.drop_rel(slide_id.rId)
    del presentation.slides._sldIdLst[index]
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _append_title(source: bytes, title: str) -> bytes:
    presentation = Presentation(BytesIO(source))
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = title
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _presentation_with_hidden_page(visible_title: str, hidden_title: str) -> bytes:
    presentation = Presentation()
    visible = presentation.slides.add_slide(presentation.slide_layouts[5])
    visible.shapes.title.text = visible_title
    hidden = presentation.slides.add_slide(presentation.slide_layouts[5])
    hidden.shapes.title.text = hidden_title
    hidden._element.set("show", "0")
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _replace_title(source: bytes, old_title: str, new_title: str) -> bytes:
    presentation = Presentation(BytesIO(source))
    slide = next(
        slide
        for slide in presentation.slides
        if slide.shapes.title.text == old_title
    )
    slide.shapes.title.text = new_title
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _install_toolchain(
    monkeypatch: pytest.MonkeyPatch,
    *,
    source_image: bytes | None = None,
    source_image_reference_index: int = 0,
) -> None:
    def convert(source: bytes, page: SourcePage) -> NormalizedPageContent:
        presentation = Presentation(BytesIO(source))
        title = presentation.slides[page.page_number - 1].shapes.title.text
        images = () if source_image is None else (
            NormalizedImage(
                source_image_reference_index,
                "可继承来源图",
                "image/png",
                "ppt/media/inherited.png",
                source_image,
            ),
        )
        return NormalizedPageContent((title,), (), (), images)

    def render(
        source: bytes, *, toolchain: object, pages: tuple[SourcePage, ...]
    ) -> tuple[StandardPageRender, ...]:
        del source, toolchain
        page_number = pages[0].page_number
        return (
            StandardPageRender(
                page_number=page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=f"render-{page_number}".encode(),
            ),
        )

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", convert)
    monkeypatch.setattr("pptextract.ingest_workflow.render_standard_pages", render)


def _upload(
    client: TestClient,
    source: bytes,
    *,
    key: str,
    document_id: str | None = None,
) -> dict[str, str]:
    route = (
        "/api/v1/documents"
        if document_id is None
        else f"/api/v1/documents/{document_id}/versions"
    )
    response = client.post(
        route,
        headers={"Idempotency-Key": key},
        files={"file": (f"{key}.pptx", source, PPTX_MEDIA_TYPE)},
    )
    assert response.status_code == 202
    return response.json()


def _pages_by_title(client: TestClient) -> dict[str, dict[str, object]]:
    pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"]
    return {str(page["title"]): page for page in pages}


def _record_historical_review(
    settings: Settings,
    *,
    page_id: str,
    status: str,
    actor_id: str,
    reviewed_at: str,
    exclusion_reason: str | None = None,
) -> tuple[str, str]:
    with transaction(settings) as connection:
        page_version = connection.execute(
            """
            SELECT page_version_id, version_id
            FROM page_versions WHERE page_id = ?
            ORDER BY created_at DESC LIMIT 1
            """,
            (page_id,),
        ).fetchone()
        assert page_version is not None
        snapshot_id = "snapshot-history"
        visual_ref = "visual-history"
        connection.execute(
            """
            INSERT INTO curation_snapshots (
                snapshot_id, page_version_id, snapshot_kind, source_snapshot_id,
                overview, created_at
            ) VALUES (?, ?, 'formal', NULL, ?, ?)
            """,
            (snapshot_id, page_version["page_version_id"], "经人工确认的结论", reviewed_at),
        )
        connection.execute(
            """
            INSERT INTO visual_objects (visual_ref, page_id, created_at)
            VALUES (?, ?, ?)
            """,
            (visual_ref, page_id, reviewed_at),
        )
        connection.execute(
            """
            INSERT INTO curation_snapshot_visuals (
                snapshot_id, visual_ref, position, source_kind, disposition,
                summary, visual_type, bounds_json, source_visual_ref, confirmed
            ) VALUES (?, ?, 0, 'capture', 'included', ?, 'chart', ?, NULL, 1)
            """,
            (snapshot_id, visual_ref, "人工确认的趋势", '[0.1,0.2,0.8,0.7]'),
        )
        connection.execute(
            """
            UPDATE page_versions
            SET review_status = ?, current_snapshot_id = ?, reviewed_by = ?,
                reviewed_at = ?, review_source_version_id = ?,
                exclusion_reason = ?
            WHERE page_version_id = ?
            """,
            (
                status,
                snapshot_id,
                actor_id,
                reviewed_at,
                page_version["version_id"],
                exclusion_reason,
                page_version["page_version_id"],
            ),
        )
        connection.execute(
            """
            INSERT INTO page_review_events (
                event_id, page_version_id, event_type, actor_id, occurred_at,
                source_version_id, source_page_version_id, snapshot_id,
                reason, note
            ) VALUES ('review-history', ?, ?, ?, ?, ?, NULL, ?, ?, NULL)
            """,
            (
                page_version["page_version_id"],
                status,
                actor_id,
                reviewed_at,
                page_version["version_id"],
                snapshot_id,
                exclusion_reason,
            ),
        )
    return snapshot_id, visual_ref


def test_unique_unchanged_pages_keep_identity_when_reordered(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(client, _presentation("甲", "乙", "丙"), key="first")
    assert run_once(settings) is True
    before = _pages_by_title(client)

    second = _upload(
        client,
        _presentation("丙", "甲", "乙"),
        key="reordered",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    after = _pages_by_title(client)

    assert second["version_id"] != first["version_id"]
    assert {title: page["page_number"] for title, page in after.items()} == {
        "丙": 1,
        "甲": 2,
        "乙": 3,
    }
    assert {
        title: (page["page_id"], page["chunk_id"]) for title, page in after.items()
    } == {
        title: (page["page_id"], page["chunk_id"]) for title, page in before.items()
    }


def test_inserted_pages_get_new_identity_and_missing_pages_restore_their_identity(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    source = _presentation("甲", "乙", "丙")
    first = _upload(client, source, key="baseline")
    assert run_once(settings) is True
    baseline = _pages_by_title(client)

    without_middle = _remove_title(source, "乙")
    _upload(
        client,
        without_middle,
        key="remove-middle",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    after_removal = _pages_by_title(client)
    assert set(after_removal) == {"甲", "丙"}
    assert client.get(f"/api/v1/pages/{baseline['乙']['page_id']}").status_code == 404

    with_new_page = _append_title(without_middle, "丁")
    _upload(
        client,
        with_new_page,
        key="append-new",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    after_insert = _pages_by_title(client)
    assert after_insert["丁"]["page_id"] not in {
        page["page_id"] for page in baseline.values()
    }
    assert after_insert["丁"]["chunk_id"] not in {
        page["chunk_id"] for page in baseline.values()
    }

    restored_source = _append_title(source, "丁")
    _upload(
        client,
        restored_source,
        key="restore-middle",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    restored = _pages_by_title(client)

    assert (restored["乙"]["page_id"], restored["乙"]["chunk_id"]) == (
        baseline["乙"]["page_id"],
        baseline["乙"]["chunk_id"],
    )
    assert (restored["丁"]["page_id"], restored["丁"]["chunk_id"]) == (
        after_insert["丁"]["page_id"],
        after_insert["丁"]["chunk_id"],
    )


@pytest.mark.parametrize(
    ("status", "exclusion_reason"),
    (("approved", None), ("excluded", "irrelevant")),
)
def test_unchanged_page_inherits_formal_annotation_and_original_review(
    system: tuple[TestClient, Settings],
    monkeypatch: pytest.MonkeyPatch,
    status: str,
    exclusion_reason: str | None,
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(client, _presentation("继承页", subject="first"), key="review-first")
    assert run_once(settings) is True
    original = _pages_by_title(client)["继承页"]
    source_snapshot_id, visual_ref = _record_historical_review(
        settings,
        page_id=str(original["page_id"]),
        status=status,
        actor_id="curator-wang",
        reviewed_at="2026-08-20T12:00:00+00:00",
        exclusion_reason=exclusion_reason,
    )

    second = _upload(
        client,
        _presentation("继承页", subject="second"),
        key="review-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    inherited = _pages_by_title(client)["继承页"]
    detail = client.get(f"/api/v1/pages/{inherited['page_id']}").json()

    assert inherited["page_id"] == original["page_id"]
    assert inherited["chunk_id"] == original["chunk_id"]
    assert inherited["review_status"] == status
    assert detail["review"] == {
        "status": status,
        "reviewed_by": "curator-wang",
        "reviewed_at": "2026-08-20T12:00:00+00:00",
        "source_version_id": first["version_id"],
        "inherited_from_page_version_id": detail["review"][
            "inherited_from_page_version_id"
        ],
        "exclusion_reason": exclusion_reason,
        "exclusion_note": None,
    }
    assert detail["review"]["inherited_from_page_version_id"] is not None
    inherited_pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "inherited"}
    )
    assert inherited_pages.status_code == 200
    inherited_payload = inherited_pages.json()["pages"]
    assert [item["page_id"] for item in inherited_payload] == [inherited["page_id"]]
    assert inherited_payload[0]["review"] == detail["review"]
    assert detail["annotation"]["source_snapshot_id"] == source_snapshot_id
    assert detail["annotation"]["overview"] == "经人工确认的结论"
    assert detail["annotation"]["visuals"] == [
        {
            "visual_ref": visual_ref,
            "position": 0,
            "source_kind": "capture",
            "disposition": "included",
            "summary": "人工确认的趋势",
            "visual_type": "chart",
            "bounds": [0.1, 0.2, 0.8, 0.7],
            "source_visual_ref": None,
            "confirmed": True,
        }
    ]
    assert detail["prefill"] is None

    with connect(settings) as connection:
        event = connection.execute(
            """
            SELECT events.event_type, events.actor_id, events.source_version_id,
                   events.source_page_version_id
            FROM page_review_events AS events
            JOIN page_versions AS pv ON pv.page_version_id = events.page_version_id
            WHERE pv.version_id = ?
            """,
            (second["version_id"],),
        ).fetchone()
    assert dict(event) == {
        "event_type": "inherited",
        "actor_id": None,
        "source_version_id": first["version_id"],
        "source_page_version_id": detail["review"]["inherited_from_page_version_id"],
    }


def test_unchanged_page_inherits_frozen_source_content_and_review_facts(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(
        client,
        _presentation("来源继承页", subject="first"),
        key="source-review-first",
    )
    assert run_once(settings) is True
    original = _pages_by_title(client)["来源继承页"]
    page_id = str(original["page_id"])

    saved = client.post(
        f"/api/v1/pages/{page_id}/curation/snapshots",
        headers={"X-Actor-ID": "curator-source"},
        json={
            "base_snapshot_id": None,
            "titles": ["人工修订的继承标题"],
            "body": [],
        },
    ).json()["curation"]["current_snapshot"]
    for route, actor in (
        ("source-confirmation", "curator-confirm"),
        ("source-review", "curator-review"),
    ):
        response = client.post(
            f"/api/v1/pages/{page_id}/curation/{route}",
            headers={"X-Actor-ID": actor},
            json={"snapshot_id": saved["snapshot_id"]},
        )
        assert response.status_code == 200
    approved = client.post(
        f"/api/v1/pages/{page_id}/approve",
        headers={"X-Actor-ID": "curator-approve"},
        json={"snapshot_id": saved["snapshot_id"]},
    )
    assert approved.status_code == 200

    _upload(
        client,
        _presentation("来源继承页", subject="second"),
        key="source-review-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    inherited = _pages_by_title(client)["来源继承页"]
    detail = client.get(f"/api/v1/pages/{inherited['page_id']}").json()
    snapshot = detail["curation"]["current_snapshot"]

    assert inherited["review_status"] == "approved"
    assert snapshot["snapshot_id"] != saved["snapshot_id"]
    assert snapshot["source_snapshot_id"] == saved["snapshot_id"]
    assert snapshot["source_content"]["titles"] == ["人工修订的继承标题"]
    assert snapshot["created_by"] == "curator-source"
    assert snapshot["source_confirmation"]["actor_id"] == "curator-confirm"
    assert snapshot["source_review"]["actor_id"] == "curator-review"
    assert detail["curation"]["chunk_body"] == {
        "nonempty": True,
        "preview": "人工修订的继承标题",
    }
    assert detail["curation"]["blockers"] == []


def test_unchanged_page_inherits_frozen_source_image_disposition_and_original_asset(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    original_bytes = b"inherited-anydoc-original"
    _install_toolchain(
        monkeypatch,
        source_image=original_bytes,
        source_image_reference_index=7,
    )
    first = _upload(
        client,
        _presentation("图片来源继承页", subject="first"),
        key="image-source-first",
    )
    assert run_once(settings) is True
    original = _pages_by_title(client)["图片来源继承页"]
    page_id = str(original["page_id"])
    state = client.post(
        f"/api/v1/pages/{page_id}/curation/snapshots",
        json={"base_snapshot_id": None, "titles": ["图片来源继承页"], "body": []},
    ).json()["curation"]
    state = client.post(
        f"/api/v1/pages/{page_id}/curation/source-confirmation",
        json={"snapshot_id": state["current_snapshot"]["snapshot_id"]},
    ).json()["curation"]
    original_source_ref = state["image_sources"]["items"][0]["source_ref"]
    state = client.post(
        f"/api/v1/pages/{page_id}/curation/image-sources/{original_source_ref}",
        headers={"X-Actor-ID": "curator-image"},
        json={
            "base_snapshot_id": state["current_snapshot"]["snapshot_id"],
            "disposition": "included",
            "summary": "公开图片展示可继承的原始资产契约。",
        },
    ).json()["curation"]
    visual_ref = state["image_sources"]["items"][0]["visual_ref"]
    asset_sha256 = state["image_sources"]["items"][0]["object_sha256"]
    state = client.post(
        f"/api/v1/pages/{page_id}/curation/source-review",
        json={"snapshot_id": state["current_snapshot"]["snapshot_id"]},
    ).json()["curation"]
    assert client.post(
        f"/api/v1/pages/{page_id}/approve",
        json={"snapshot_id": state["current_snapshot"]["snapshot_id"]},
    ).status_code == 200

    _upload(
        client,
        _presentation("图片来源继承页", subject="second"),
        key="image-source-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    inherited = _pages_by_title(client)["图片来源继承页"]
    detail = client.get(f"/api/v1/pages/{inherited['page_id']}").json()
    item = detail["curation"]["image_sources"]["items"][0]

    assert inherited["review_status"] == "approved"
    assert item["source_ref"] != original_source_ref
    assert item["disposition"] == "included"
    assert item["summary"] == "公开图片展示可继承的原始资产契约。"
    assert item["visual_ref"] == visual_ref
    assert item["object_sha256"] == asset_sha256
    assert detail["curation"]["blockers"] == []
    assert detail["annotation"]["visuals"][0]["source_image_ref"] == item["source_ref"]
    assert detail["annotation"]["visuals"][0]["asset"] == {
        "sha256": asset_sha256,
        "media_type": "image/png",
        "size_bytes": len(original_bytes),
        "byte_contract": "anydoc_original",
    }


def test_changed_page_keeps_identity_but_only_receives_unconfirmed_prefill(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(client, _presentation("变化前"), key="changed-first")
    assert run_once(settings) is True
    original = _pages_by_title(client)["变化前"]
    source_snapshot_id, old_visual_ref = _record_historical_review(
        settings,
        page_id=str(original["page_id"]),
        status="approved",
        actor_id="curator-chen",
        reviewed_at="2026-08-20T13:00:00+00:00",
    )

    _upload(
        client,
        _presentation("变化后"),
        key="changed-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    changed = _pages_by_title(client)["变化后"]
    detail = client.get(f"/api/v1/pages/{changed['page_id']}").json()

    assert changed["page_id"] == original["page_id"]
    assert changed["chunk_id"] == original["chunk_id"]
    assert changed["review_status"] == "pending"
    assert detail["review"] == {
        "status": "pending",
        "reviewed_by": None,
        "reviewed_at": None,
        "source_version_id": None,
        "inherited_from_page_version_id": None,
        "exclusion_reason": None,
        "exclusion_note": None,
    }
    assert detail["annotation"] is None
    assert detail["prefill"]["kind"] == "prefill"
    assert detail["prefill"]["source_snapshot_id"] == source_snapshot_id
    assert detail["prefill"]["overview"] == "经人工确认的结论"
    (visual,) = detail["prefill"]["visuals"]
    assert visual["visual_ref"] != old_visual_ref
    assert visual["source_visual_ref"] == old_visual_ref
    assert visual["bounds"] == [0.1, 0.2, 0.8, 0.7]
    assert visual["confirmed"] is False


def test_rollback_to_historical_content_uses_the_latest_matching_review_history(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(client, _presentation("历史内容"), key="history-first")
    assert run_once(settings) is True
    original = _pages_by_title(client)["历史内容"]
    _record_historical_review(
        settings,
        page_id=str(original["page_id"]),
        status="approved",
        actor_id="curator-li",
        reviewed_at="2026-08-20T14:00:00+00:00",
    )

    changed = _upload(
        client,
        _presentation("当前变化内容"),
        key="history-changed",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    assert _pages_by_title(client)["当前变化内容"]["review_status"] == "pending"

    rollback = client.post(
        f"/api/v1/documents/{first['document_id']}"
        f"/versions/{first['version_id']}/rollback",
        headers={"X-Actor-ID": "operator-li", "Idempotency-Key": "rollback-history"},
        json={"reason": "恢复已经审核的内容"},
    )
    assert rollback.status_code == 202
    assert rollback.json()["version_id"] not in {
        first["version_id"],
        changed["version_id"],
    }
    assert run_once(settings) is True

    restored = _pages_by_title(client)["历史内容"]
    detail = client.get(f"/api/v1/pages/{restored['page_id']}").json()
    assert restored["page_id"] == original["page_id"]
    assert restored["review_status"] == "approved"
    assert detail["review"]["reviewed_by"] == "curator-li"
    assert detail["review"]["source_version_id"] == first["version_id"]
    assert detail["annotation"]["overview"] == "经人工确认的结论"


@pytest.mark.product_fault
def test_activation_failure_rolls_back_page_versions_inheritance_and_current_pointer(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first = _upload(
        client,
        _presentation("事务页", subject="baseline"),
        key="atomic-first",
    )
    assert run_once(settings) is True
    original = _pages_by_title(client)["事务页"]
    _record_historical_review(
        settings,
        page_id=str(original["page_id"]),
        status="approved",
        actor_id="curator-zhao",
        reviewed_at="2026-08-20T15:00:00+00:00",
    )
    second = _upload(
        client,
        _presentation("事务页", subject="replacement"),
        key="atomic-second",
        document_id=first["document_id"],
    )
    with transaction(settings) as connection:
        connection.execute(
            """
            CREATE TRIGGER reject_inheritance_event
            BEFORE INSERT ON page_review_events
            WHEN NEW.event_type = 'inherited'
            BEGIN
                SELECT RAISE(ABORT, 'injected activation failure');
            END
            """
        )

    assert run_once(settings) is True

    document = client.get(f"/api/v1/documents/{first['document_id']}").json()
    failed_version = client.get(
        f"/api/v1/documents/{first['document_id']}"
        f"/versions/{second['version_id']}"
    ).json()
    visible = _pages_by_title(client)["事务页"]
    assert document["current_version_id"] == first["version_id"]
    assert failed_version["status"] == "failed"
    assert visible["page_id"] == original["page_id"]
    assert visible["review_status"] == "approved"
    with connect(settings) as connection:
        assert connection.execute(
            "SELECT COUNT(*) FROM page_versions WHERE version_id = ?",
            (second["version_id"],),
        ).fetchone()[0] == 0
        assert connection.execute(
            """
            SELECT COUNT(*) FROM page_review_events AS events
            JOIN page_versions AS pv ON pv.page_version_id = events.page_version_id
            WHERE pv.version_id = ?
            """,
            (second["version_id"],),
        ).fetchone()[0] == 0


def test_changed_hidden_page_restores_identity_and_prefill_when_enabled(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    source = _presentation_with_hidden_page("可见页", "隐藏页旧内容")
    first = _upload(client, source, key="hidden-history-first")
    assert run_once(settings) is True
    first_enable = client.post(
        f"/api/v1/documents/{first['document_id']}"
        f"/versions/{first['version_id']}/source-pages/2/enable",
        headers={"X-Actor-ID": "curator-sun", "Idempotency-Key": "enable-first-hidden"},
    )
    assert first_enable.status_code == 202
    assert run_once(settings) is True
    original = _pages_by_title(client)["隐藏页旧内容"]
    source_snapshot_id, old_visual_ref = _record_historical_review(
        settings,
        page_id=str(original["page_id"]),
        status="approved",
        actor_id="curator-sun",
        reviewed_at="2026-08-20T16:00:00+00:00",
    )

    second_source = _replace_title(source, "隐藏页旧内容", "隐藏页新内容")
    second = _upload(
        client,
        second_source,
        key="hidden-history-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    second_enable = client.post(
        f"/api/v1/documents/{first['document_id']}"
        f"/versions/{second['version_id']}/source-pages/2/enable",
        headers={"X-Actor-ID": "curator-sun", "Idempotency-Key": "enable-second-hidden"},
    )
    assert second_enable.status_code == 202
    assert run_once(settings) is True

    changed = _pages_by_title(client)["隐藏页新内容"]
    detail = client.get(f"/api/v1/pages/{changed['page_id']}").json()
    assert changed["page_id"] == original["page_id"]
    assert changed["chunk_id"] == original["chunk_id"]
    assert changed["review_status"] == "pending"
    assert detail["annotation"] is None
    assert detail["prefill"]["source_snapshot_id"] == source_snapshot_id
    (visual,) = detail["prefill"]["visuals"]
    assert visual["visual_ref"] != old_visual_ref
    assert visual["source_visual_ref"] == old_visual_ref
    assert visual["confirmed"] is False
