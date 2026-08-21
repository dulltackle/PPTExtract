from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptextract.ingestion import convert_presentation
from pptextract.pptx_projection import list_source_pages
from tests.support.synthetic_pptx import build_public_contract_presentation


def test_public_synthetic_fixture_covers_required_source_scenarios() -> None:
    source = build_public_contract_presentation()
    presentation = Presentation(BytesIO(source))

    assert len(presentation.slides) == 9
    assert any(shape.has_chart for shape in presentation.slides[3].shapes)
    assert any(shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in presentation.slides[3].shapes)
    assert presentation.slides[4].shapes[-1].has_table
    assert presentation.slides[5].notes_slide.notes_text_frame.text == (
        "仅供契约测试的公开演讲者备注"
    )
    assert list_source_pages(source)[6].hidden is True
    assert presentation.slides[8].shapes[-1].text_frame.paragraphs[0].runs[0].font.name == (
        "PPTExtract Missing Contract Font"
    )


def test_public_fixture_locks_page_isolation_duplicates_and_reordering() -> None:
    original = convert_presentation(build_public_contract_presentation())
    reordered = convert_presentation(
        build_public_contract_presentation(order=(8, 1, 2, 3, 4, 5, 6, 7, 9))
    )

    assert [image.alt_text for image in original[1].content.images] == [
        "甲页第一处图片",
        "甲页第二处图片",
    ]
    assert [image.alt_text for image in original[2].content.images] == ["乙页图片"]
    assert original[1].content.body == ("甲页唯一正文",)
    assert original[2].content.body == ("乙页唯一正文",)
    assert original[0].content.body == ("纯文字页唯一正文", "公开项目一\n公开项目二")
    assert original[0].fingerprint == original[7].fingerprint
    assert [page.source.source_slide_id for page in reordered[:2]] == [263, 256]
    assert reordered[0].fingerprint == original[7].fingerprint
