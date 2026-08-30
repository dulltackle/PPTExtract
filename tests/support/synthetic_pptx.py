from __future__ import annotations

from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


def build_minimal_presentation() -> bytes:
    """构造只含虚构内容的最小公开 PPTX 契约夹具。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "公开合成页一"

    hidden = presentation.slides.add_slide(presentation.slide_layouts[5])
    hidden.shapes.title.text = "公开合成隐藏页"
    hidden._element.set("show", "0")

    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def build_plain_text_presentation(
    *,
    title: str = "公开首次摄取",
    body_text: str = "这是可公开验证的单页纯文字内容。",
) -> bytes:
    """构造 Issue #20 首次摄取快乐路径使用的单页纯文字 PPTX。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = title
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(1))
    body.text = body_text

    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def build_repeated_footer_presentation() -> bytes:
    """构造三页公开合成内容，用于人工确认重复页脚噪声。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for page_number in range(1, 4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"公开页 {page_number}"
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(1))
        body.text = f"第 {page_number} 页独有正文"
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7), Inches(8), Inches(0.3))
        footer.text = "公开合成重复页脚"

    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def build_rendering_warning_presentation() -> bytes:
    """构造同时含缺失字体与动画时间线的公开渲染警告夹具。"""
    from xml.etree import ElementTree
    from zipfile import ZIP_DEFLATED, ZipFile

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    font_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    font_slide.shapes.title.text = "公开缺失字体页"
    text_box = font_slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11), Inches(1))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "此文字故意引用不存在的字体"
    run.font.name = "PPTExtract Missing Contract Font"

    animation_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    animation_slide.shapes.title.text = "公开动画扁平化页"

    stream = BytesIO()
    presentation.save(stream)
    output = BytesIO()
    presentation_namespace = "http://schemas.openxmlformats.org/presentationml/2006/main"
    with ZipFile(BytesIO(stream.getvalue())) as package, ZipFile(
        output, "w", compression=ZIP_DEFLATED
    ) as rewritten:
        for entry in package.infolist():
            content = package.read(entry.filename)
            if entry.filename == "ppt/slides/slide2.xml":
                slide = ElementTree.fromstring(content)
                timing = ElementTree.SubElement(
                    slide, f"{{{presentation_namespace}}}timing"
                )
                ElementTree.SubElement(timing, f"{{{presentation_namespace}}}tnLst")
                content = ElementTree.tostring(slide, encoding="utf-8", xml_declaration=True)
            rewritten.writestr(entry, content)
    return output.getvalue()


def build_table_font_presentation() -> bytes:
    """构造表格单元格显式缺失字体的公开契约夹具。"""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(5), Inches(1)).table
    run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
    run.text = "公开表格字体"
    run.font.name = "PPTExtract Missing Table Font"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def build_installed_font_glyph_fallback_presentation() -> bytes:
    """构造字体家族已安装、但中文字形会由渲染器回退的公开夹具。"""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "ABC 公开中文字形回退"
    run.font.name = "Liberation Sans"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def build_conversion_presentation() -> tuple[bytes, bytes]:
    """构造带重复图片引用的公开转换契约夹具。"""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "公开图片页"
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5), Inches(0.6))
    body.text = "第一段正文"

    image_stream = BytesIO()
    Image.new("RGB", (24, 16), (224, 48, 64)).save(image_stream, format="PNG")
    image_bytes = image_stream.getvalue()
    first = slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.5), Inches(2.2))
    first._element.nvPicPr.cNvPr.set("descr", "第一处图片引用")
    second = slide.shapes.add_picture(BytesIO(image_bytes), Inches(3), Inches(2.2))
    second._element.nvPicPr.cNvPr.set("descr", "第二处图片引用")

    table = slide.shapes.add_table(3, 3, Inches(6), Inches(1.4), Inches(6), Inches(2)).table
    table.first_row = True
    table.cell(0, 0).text = "合并表头"
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 2).text = "指标"
    table.cell(1, 0).text = "甲"
    table.cell(1, 1).text = "10"
    table.cell(1, 2).text = "20"
    table.cell(2, 0).text = "乙"
    table.cell(2, 1).text = "30"
    table.cell(2, 2).text = "40"
    slide.notes_slide.notes_text_frame.text = "公开演讲者备注"

    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue(), image_bytes


def build_image_curation_presentation() -> bytes:
    """构造图片来源真实浏览器处置使用的三页公开 PPTX。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    image_payloads: list[bytes] = []
    for color in ((32, 112, 208), (82, 156, 111), (224, 96, 64)):
        stream = BytesIO()
        Image.new("RGB", (64, 40), color).save(stream, format="PNG")
        image_payloads.append(stream.getvalue())

    keep_page = _add_title_and_body(presentation, "公开单项保留页", "保留流程正文")
    _add_picture(keep_page, image_payloads[0], "单项保留来源图", 0.7)

    ignore_page = _add_title_and_body(presentation, "公开单项忽略页", "忽略流程正文")
    _add_picture(ignore_page, image_payloads[1], "单项忽略来源图", 0.7)

    mixed_page = _add_title_and_body(presentation, "公开混合处置页", "混合流程正文")
    _add_picture(mixed_page, image_payloads[2], "重复对象第一处引用", 0.7)
    _add_picture(mixed_page, image_payloads[2], "重复对象第二处引用", 4.1)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def build_product_acceptance_presentation(
    *, order: tuple[int, ...] | None = None
) -> bytes:
    """构造产品级黑盒门禁的单文档四页合成 PPTX。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    _add_title_and_body(presentation, "公开重复结论页", "相同内容由人分别批准或排除。")

    image_stream = BytesIO()
    Image.new("RGB", (80, 48), (36, 112, 176)).save(image_stream, format="PNG")
    image_page = _add_title_and_body(
        presentation,
        "公开来源图片页",
        "此页验证 AnyDoc 来源图片逐项处置。",
    )
    _add_picture(image_page, image_stream.getvalue(), "公开蓝色来源图片", 0.8)

    _add_title_and_body(
        presentation,
        "公开人工框选页",
        "此页需要从标准页渲染结果补充视觉结论。",
    )

    _add_title_and_body(presentation, "公开重复结论页", "相同内容由人分别批准或排除。")

    output = BytesIO()
    presentation.save(output)
    result = output.getvalue()
    return result if order is None else _reorder_slides(result, order)


def build_public_contract_presentation(
    *, order: tuple[int, ...] | None = None
) -> bytes:
    """生成 Issue #19 的公开、虚构、可重复契约夹具。"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    text_page = _add_title_and_body(presentation, "公开纯文字页", "纯文字页唯一正文")
    _add_bullet_list(text_page, ("公开项目一", "公开项目二"))

    shared_image = BytesIO()
    Image.new("RGB", (64, 40), (32, 112, 208)).save(shared_image, format="PNG")
    image_bytes = shared_image.getvalue()
    first_image_page = _add_title_and_body(presentation, "公开图片页甲", "甲页唯一正文")
    _add_picture(first_image_page, image_bytes, "甲页第一处图片", 0.7)
    _add_picture(first_image_page, image_bytes, "甲页第二处图片", 4.1)
    second_image_page = _add_title_and_body(presentation, "公开图片页乙", "乙页唯一正文")
    _add_picture(second_image_page, image_bytes, "乙页图片", 0.7)

    visual_page = _add_title_and_body(presentation, "公开图表组合页", "图表和组合形状")
    chart_data = ChartData()
    chart_data.categories = ["一月", "二月", "三月"]
    chart_data.add_series("公开序列", (2, 5, 3))
    visual_page.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.7),
        Inches(2),
        Inches(5.5),
        Inches(4),
        chart_data,
    )
    group = visual_page.shapes.add_group_shape()
    left = group.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7),
        Inches(2.2),
        Inches(2),
        Inches(1.2),
    )
    left.text = "组合甲"
    arrow = group.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(9.2),
        Inches(2.45),
        Inches(1.4),
        Inches(0.7),
    )
    arrow.text = ""
    right = group.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.8),
        Inches(2.2),
        Inches(2),
        Inches(1.2),
    )
    right.text = "组合乙"

    table_page = _add_title_and_body(presentation, "公开复杂表格页", "合并网格")
    table = table_page.shapes.add_table(
        4, 4, Inches(0.8), Inches(2), Inches(11.7), Inches(3.5)
    ).table
    table.first_row = True
    table.cell(0, 0).text = "双列表头"
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 2).text = "纵向表头"
    table.cell(0, 2).merge(table.cell(1, 2))
    table.cell(0, 3).text = "指标"
    for row in range(1, 4):
        for column in range(4):
            if not table.cell(row, column).is_spanned:
                table.cell(row, column).text = f"R{row}C{column}"

    notes_page = _add_title_and_body(presentation, "公开备注页", "备注页正文")
    notes_page.notes_slide.notes_text_frame.text = "仅供契约测试的公开演讲者备注"

    hidden_page = _add_title_and_body(presentation, "公开隐藏页", "隐藏页正文")
    hidden_page._element.set("show", "0")

    duplicate_page = _add_title_and_body(presentation, "公开纯文字页", "纯文字页唯一正文")
    _add_bullet_list(duplicate_page, ("公开项目一", "公开项目二"))
    assert text_page.shapes.title.text == duplicate_page.shapes.title.text

    missing_font_page = presentation.slides.add_slide(presentation.slide_layouts[5])
    missing_font_page.shapes.title.text = "公开缺失字体页"
    text_box = missing_font_page.shapes.add_textbox(
        Inches(0.8), Inches(2), Inches(11), Inches(1)
    )
    paragraph = text_box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "此文字故意引用不存在的字体"
    run.font.name = "PPTExtract Missing Contract Font"

    stream = BytesIO()
    presentation.save(stream)
    result = stream.getvalue()
    if order is not None:
        result = _reorder_slides(result, order)
    return result


def _add_title_and_body(presentation: Presentation, title: str, body: str):  # type: ignore[no-untyped-def]
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = title
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(7), Inches(0.6))
    text_box.text = body
    return slide


def _add_picture(slide, image_bytes: bytes, alt_text: str, left_inches: float) -> None:  # type: ignore[no-untyped-def]
    picture = slide.shapes.add_picture(
        BytesIO(image_bytes), Inches(left_inches), Inches(2.3), width=Inches(2.8)
    )
    picture._element.nvPicPr.cNvPr.set("descr", alt_text)


def _add_bullet_list(slide, items: tuple[str, ...]) -> None:  # type: ignore[no-untyped-def]
    text_box = slide.shapes.add_textbox(Inches(0.9), Inches(2), Inches(5), Inches(1.8))
    for index, item in enumerate(items):
        paragraph = (
            text_box.text_frame.paragraphs[0]
            if index == 0
            else text_box.text_frame.add_paragraph()
        )
        paragraph.text = item
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        paragraph._p.get_or_add_pPr().insert(0, bullet)


def _reorder_slides(pptx_bytes: bytes, order: tuple[int, ...]) -> bytes:
    from copy import deepcopy
    from xml.etree import ElementTree
    from zipfile import ZIP_DEFLATED, ZipFile

    namespace = "http://schemas.openxmlformats.org/presentationml/2006/main"
    with ZipFile(BytesIO(pptx_bytes)) as source:
        presentation = ElementTree.fromstring(source.read("ppt/presentation.xml"))
        slide_id_list = presentation.find(f"{{{namespace}}}sldIdLst")
        if slide_id_list is None or sorted(order) != list(range(1, len(slide_id_list) + 1)):
            raise ValueError("页序必须完整且不重复")
        originals = tuple(slide_id_list)
        slide_id_list[:] = [deepcopy(originals[index - 1]) for index in order]
        output = BytesIO()
        with ZipFile(output, "w", compression=ZIP_DEFLATED) as target:
            for entry in source.infolist():
                content = source.read(entry.filename)
                if entry.filename == "ppt/presentation.xml":
                    content = ElementTree.tostring(
                        presentation, encoding="utf-8", xml_declaration=True
                    )
                target.writestr(entry, content)
    return output.getvalue()
