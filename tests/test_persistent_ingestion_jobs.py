from __future__ import annotations

from collections.abc import Iterator
from dataclasses import replace
from io import BytesIO
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from pptextract.api import create_app
from pptextract.config import Settings
from pptextract.conversion import NormalizedPageContent
from pptextract.db import connect, initialize_database, transaction
from pptextract.jobs import claim_next_job, enqueue_job, finish_job
from pptextract.object_store import LocalObjectStore
from pptextract.pptx_projection import SourcePage
from pptextract.rendering import StandardPageRender
from pptextract.worker import run_once
from tests.support.synthetic_pptx import build_minimal_presentation, build_plain_text_presentation

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture
def system(tmp_path: Path) -> Iterator[tuple[TestClient, Settings]]:
    settings = replace(Settings.for_test(tmp_path), job_retry_base_seconds=0)
    with TestClient(create_app(settings)) as client:
        yield client, settings


def _accept_ingestion(client: TestClient) -> dict[str, str]:
    response = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "persistent-job"},
        files={
            "file": (
                "persistent-job.pptx",
                build_plain_text_presentation(),
                PPTX_MEDIA_TYPE,
            )
        },
    )
    assert response.status_code == 202
    return response.json()


def _build_two_page_presentation() -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for page_number in (1, 2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"断点页 {page_number}"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def test_job_query_exposes_six_states_independently_from_version_state(
    system: tuple[TestClient, Settings],
) -> None:
    client, settings = system
    accepted = _accept_ingestion(client)
    job_id = accepted["job_id"]

    assert client.get(f"/api/v1/jobs/{job_id}").json()["status"] == "queued"

    for status in ("running", "requires_action", "succeeded", "failed", "cancelled"):
        with transaction(settings) as connection:
            connection.execute(
                "UPDATE jobs SET status = ? WHERE job_id = ?",
                (status, job_id),
            )

        task = client.get(f"/api/v1/jobs/{job_id}")
        version = client.get(
            f"/api/v1/documents/{accepted['document_id']}"
            f"/versions/{accepted['version_id']}"
        )

        assert task.status_code == 200
        assert task.json()["status"] == status
        assert version.json()["status"] == "processing"


def test_transient_conversion_failure_is_retried_in_the_same_observable_job(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = _accept_ingestion(client)

    def unavailable_conversion(*_args: object, **_kwargs: object) -> object:
        raise OSError("temporary converter outage")

    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        unavailable_conversion,
    )

    assert run_once(settings) is True

    task = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    version = client.get(
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}"
    ).json()
    assert task["status"] == "queued"
    assert task["attempts"] == 1
    assert task["next_retry_at"] is not None
    assert task["error"] == {
        "attempt": 1,
        "code": "conversion_temporarily_unavailable",
        "message": "第 1 页转换暂时失败。",
        "page_number": 1,
        "phase": "conversion",
        "retryable": True,
    }
    assert task["progress"] == {
        "phase": "conversion",
        "completed_pages": 0,
        "total_pages": 1,
        "pages": [
            {
                "page_number": 1,
                "phase": "source_manifest",
                "status": "completed",
            }
        ],
    }
    assert version["status"] == "processing"


def test_retry_resumes_page_checkpoints_without_exposing_or_duplicating_results(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    response = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "resume-checkpoints"},
        files={"file": ("two-pages.pptx", _build_two_page_presentation(), PPTX_MEDIA_TYPE)},
    )
    accepted = response.json()
    conversion_calls: list[int] = []
    rendering_calls: list[int] = []

    def convert(_source: bytes, page: SourcePage) -> NormalizedPageContent:
        page_number = page.page_number
        conversion_calls.append(page_number)
        return NormalizedPageContent(
            titles=(f"断点页 {page_number}",),
            body=(),
            tables=(),
            images=(),
        )

    def render(
        _source: bytes, *, toolchain: object, pages: tuple[SourcePage, ...]
    ) -> tuple[StandardPageRender, ...]:
        del toolchain
        page_number = pages[0].page_number
        rendering_calls.append(page_number)
        if page_number == 2 and rendering_calls.count(2) == 1:
            raise OSError("temporary renderer outage")
        return (
            StandardPageRender(
                page_number=page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=f"render-{page_number}".encode(),
            ),
        )

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", convert)
    monkeypatch.setattr("pptextract.ingest_workflow.render_standard_pages", render)

    assert run_once(settings) is True
    document_during_retry = client.get(
        f"/api/v1/documents/{accepted['document_id']}"
    ).json()
    pages_during_retry = client.get("/api/v1/curation/pages").json()["pages"]
    assert document_during_retry["current_version_id"] is None
    assert pages_during_retry == []
    task_during_retry = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    assert task_during_retry["progress"]["pages"] == [
        {"page_number": 1, "phase": "rendering", "status": "completed"},
        {"page_number": 2, "phase": "conversion", "status": "completed"},
    ]

    assert run_once(settings) is True

    task = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    document = client.get(f"/api/v1/documents/{accepted['document_id']}").json()
    pages = client.get("/api/v1/curation/pages").json()["pages"]
    assert task["status"] == "succeeded"
    assert task["attempts"] == 2
    assert conversion_calls == [1, 2]
    assert rendering_calls == [1, 2, 2]
    assert document["current_version_id"] == accepted["version_id"]
    assert [page["page_number"] for page in pages] == [1, 2]
    assert len({page["page_id"] for page in pages}) == 2
    assert len({page["chunk_id"] for page in pages}) == 2


def test_expired_lease_takeover_rejects_updates_from_the_stale_claim(tmp_path: Path) -> None:
    first_worker = replace(Settings.for_test(tmp_path), worker_id="worker-one")
    second_worker = replace(first_worker, worker_id="worker-two")
    initialize_database(first_worker)
    job_id = enqueue_job(
        first_worker,
        kind="system.noop",
        payload={},
        actor_id="operator",
        idempotency_key="lease-takeover",
    )
    first_claim = claim_next_job(first_worker)
    assert first_claim is not None
    with transaction(first_worker) as connection:
        connection.execute(
            "UPDATE jobs SET lease_expires_at = '2000-01-01T00:00:00+00:00' WHERE job_id = ?",
            (job_id,),
        )

    with pytest.raises(RuntimeError, match="未持有"):
        finish_job(first_worker, first_claim, succeeded=True)

    second_claim = claim_next_job(second_worker)
    assert second_claim is not None
    assert second_claim.job_id == first_claim.job_id
    assert second_claim.lease_token != first_claim.lease_token

    finish_job(second_worker, second_claim, succeeded=True)

    with transaction(first_worker) as connection:
        row = connection.execute(
            "SELECT status, attempts FROM jobs WHERE job_id = ?", (job_id,)
        ).fetchone()
    assert dict(row) == {"status": "succeeded", "attempts": 2}


def test_expired_lease_at_max_attempts_fails_the_job_and_version(
    system: tuple[TestClient, Settings],
) -> None:
    client, settings = system
    accepted = _accept_ingestion(client)

    for expected_attempt in (1, 2, 3):
        claim = claim_next_job(settings)
        assert claim is not None
        assert claim.attempts == expected_attempt
        with transaction(settings) as connection:
            connection.execute(
                """
                UPDATE jobs
                SET lease_expires_at = '2000-01-01T00:00:00+00:00'
                WHERE job_id = ?
                """,
                (accepted["job_id"],),
            )

    assert claim_next_job(settings) is None

    task = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    version = client.get(
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}"
    ).json()
    assert task["status"] == "failed"
    assert task["attempts"] == 3
    assert task["error"] == {
        "attempt": 3,
        "code": "lease_expired",
        "message": "任务租约过期且已耗尽尝试次数。",
        "phase": "source_manifest",
        "retryable": True,
    }
    assert version["status"] == "failed"


def test_transient_failure_exhausts_the_bounded_attempts_and_fails_the_version(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = _accept_ingestion(client)

    def unavailable_conversion(*_args: object, **_kwargs: object) -> object:
        raise OSError("converter remains unavailable")

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", unavailable_conversion)

    assert [run_once(settings) for _ in range(3)] == [True, True, True]

    task = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    version = client.get(
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}"
    ).json()
    assert task["status"] == "failed"
    assert task["attempts"] == 3
    assert task["next_retry_at"] is None
    assert task["error"]["attempt"] == 3
    assert task["error"]["retryable"] is True
    assert version["status"] == "failed"


def test_transient_storage_failure_retries_while_permanent_conversion_failure_does_not(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = _accept_ingestion(client)

    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent((), (), (), ()),
    )
    monkeypatch.setattr(
        "pptextract.ingest_workflow.render_standard_pages",
        lambda *_args, pages, **_kwargs: (
            StandardPageRender(
                page_number=pages[0].page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=b"render",
            ),
        ),
    )
    original_put = LocalObjectStore.put
    put_attempts = 0

    def flaky_put(store: LocalObjectStore, payload: bytes):  # type: ignore[no-untyped-def]
        nonlocal put_attempts
        put_attempts += 1
        if put_attempts == 1:
            raise OSError("temporary storage outage")
        return original_put(store, payload)

    monkeypatch.setattr(LocalObjectStore, "put", flaky_put)

    assert run_once(settings) is True
    retrying = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    assert retrying["status"] == "queued"
    assert retrying["error"]["phase"] == "storage"
    assert retrying["error"]["retryable"] is True
    assert run_once(settings) is True
    assert client.get(f"/api/v1/jobs/{accepted['job_id']}").json()["status"] == "succeeded"

    second = client.post(
        f"/api/v1/documents/{accepted['document_id']}/versions",
        headers={"Idempotency-Key": "permanent-conversion"},
        files={
            "file": (
                "permanent.pptx",
                build_plain_text_presentation(title="永久转换错误"),
                PPTX_MEDIA_TYPE,
            )
        },
    ).json()

    def invalid_conversion(*_args: object, **_kwargs: object) -> object:
        raise ValueError("invalid normalized content")

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", invalid_conversion)
    assert run_once(settings) is True
    failed = client.get(f"/api/v1/jobs/{second['job_id']}").json()
    assert failed["status"] == "failed"
    assert failed["attempts"] == 1
    assert failed["error"]["phase"] == "conversion"
    assert failed["error"]["retryable"] is False


def test_exhausted_retry_keeps_the_old_current_version_available(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    baseline = _accept_ingestion(client)
    with transaction(settings) as connection:
        connection.execute(
            "UPDATE jobs SET status = 'succeeded' WHERE job_id = ?",
            (baseline["job_id"],),
        )
        connection.execute(
            "UPDATE document_versions SET status = 'ready' WHERE version_id = ?",
            (baseline["version_id"],),
        )
        connection.execute(
            "UPDATE documents SET current_version_id = ? WHERE document_id = ?",
            (baseline["version_id"], baseline["document_id"]),
        )

    replacement = client.post(
        f"/api/v1/documents/{baseline['document_id']}/versions",
        headers={"Idempotency-Key": "replacement-that-fails"},
        files={
            "file": (
                "replacement.pptx",
                build_plain_text_presentation(title="失败的新版本"),
                PPTX_MEDIA_TYPE,
            )
        },
    ).json()

    def unavailable_conversion(*_args: object, **_kwargs: object) -> object:
        raise OSError("converter remains unavailable")

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", unavailable_conversion)
    assert [run_once(settings) for _ in range(3)] == [True, True, True]

    document = client.get(f"/api/v1/documents/{baseline['document_id']}").json()
    failed_version = client.get(
        f"/api/v1/documents/{baseline['document_id']}"
        f"/versions/{replacement['version_id']}"
    ).json()
    assert document["current_version_id"] == baseline["version_id"]
    assert failed_version["status"] == "failed"


def test_hidden_page_is_reported_as_skipped_without_blocking_activation(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "hidden-page-progress"},
        files={"file": ("hidden.pptx", build_minimal_presentation(), PPTX_MEDIA_TYPE)},
    ).json()
    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent((), (), (), ()),
    )
    monkeypatch.setattr(
        "pptextract.ingest_workflow.render_standard_pages",
        lambda *_args, pages, **_kwargs: (
            StandardPageRender(
                page_number=pages[0].page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=b"visible-render",
            ),
        ),
    )

    assert run_once(settings) is True

    task = client.get(f"/api/v1/jobs/{accepted['job_id']}").json()
    pages = client.get("/api/v1/curation/pages").json()["pages"]
    assert task["status"] == "succeeded"
    assert task["progress"]["pages"] == [
        {"page_number": 1, "phase": "page_fingerprint", "status": "completed"},
        {"page_number": 2, "phase": "source_manifest", "status": "skipped"},
    ]
    assert [page["page_number"] for page in pages] == [1]


def test_hidden_page_can_be_enabled_once_through_the_public_api(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "hidden-page-enable-source"},
        files={"file": ("hidden.pptx", build_minimal_presentation(), PPTX_MEDIA_TYPE)},
    ).json()
    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent(
            ("公开合成页",), ("公开正文",), (), ()
        ),
    )
    monkeypatch.setattr(
        "pptextract.ingest_workflow.render_standard_pages",
        lambda *_args, pages, **_kwargs: (
            StandardPageRender(
                page_number=pages[0].page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=f"render-{pages[0].page_number}".encode(),
            ),
        ),
    )
    assert run_once(settings) is True

    all_pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"]
    hidden = all_pages[1]
    assert hidden == {
        "page_id": None,
        "chunk_id": None,
        "document_id": accepted["document_id"],
        "version_id": accepted["version_id"],
        "page_number": 2,
        "review_status": None,
        "title": None,
        "hidden": True,
        "enabled": False,
        "source_reference": {
            "slide_id": hidden["source_reference"]["slide_id"],
            "relationship_id": hidden["source_reference"]["relationship_id"],
            "part": "ppt/slides/slide2.xml",
        },
        "enablement": {"status": "not_started", "job_id": None, "error": None},
        "rendering_warnings": {
            "total": 0,
            "pages": 0,
            "unconfirmed": 0,
            "unconfirmed_pages": 0,
        },
        "version_rendering_warnings": {
            "total": 1,
            "pages": 1,
            "unconfirmed": 1,
            "unconfirmed_pages": 1,
        },
    }
    with connect(settings) as connection:
        unprocessed = connection.execute(
            """
            SELECT source_content_json, fingerprint_sha256, render_sha256
            FROM ingestion_page_results
            WHERE version_id = ? AND page_number = 2
            """,
            (accepted["version_id"],),
        ).fetchone()
    assert tuple(unprocessed) == (None, None, None)

    path = (
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}/source-pages/2/enable"
    )
    first = client.post(
        path,
        headers={"X-Actor-ID": "curator-one", "Idempotency-Key": "enable-hidden-page"},
    )
    replay = client.post(
        path,
        headers={"X-Actor-ID": "curator-one", "Idempotency-Key": "enable-hidden-page"},
    )
    other_session = client.post(
        path,
        headers={"X-Actor-ID": "curator-two", "Idempotency-Key": "other-session"},
    )

    assert first.status_code == replay.status_code == other_session.status_code == 202
    assert replay.json() == first.json()
    assert other_session.json()["status"] == "coalesced"
    assert other_session.json()["job_id"] == first.json()["job_id"]
    queued_hidden = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"][1]
    assert queued_hidden["enablement"] == {
        "status": "queued",
        "job_id": first.json()["job_id"],
        "error": None,
    }
    queued_task = client.get(f"/api/v1/jobs/{first.json()['job_id']}").json()
    assert queued_task["progress"] == {
        "phase": "queued",
        "completed_pages": 0,
        "total_pages": 1,
        "pages": [{"page_number": 2, "phase": "queued", "status": "pending"}],
    }

    assert run_once(settings) is True

    pending_pages = client.get("/api/v1/curation/pages").json()["pages"]
    assert [page["page_number"] for page in pending_pages] == [1, 2]
    enabled_page = pending_pages[1]
    assert enabled_page["hidden"] is True
    assert enabled_page["enabled"] is True
    assert enabled_page["review_status"] == "pending"
    detail = client.get(f"/api/v1/pages/{enabled_page['page_id']}")
    assert detail.status_code == 200
    assert detail.json()["page_number"] == 2
    assert detail.json()["source_content"]["titles"] == ["公开合成页"]

    already_enabled = client.post(
        path,
        headers={"X-Actor-ID": "curator-two", "Idempotency-Key": "already-enabled"},
    )
    assert already_enabled.status_code == 200
    assert already_enabled.json()["status"] == "no_change"
    assert already_enabled.json()["page_id"] == enabled_page["page_id"]
    with connect(settings) as connection:
        counts = connection.execute(
            """
            SELECT
              (SELECT COUNT(*) FROM pages WHERE document_id = ?) AS pages,
              (SELECT COUNT(*) FROM page_versions WHERE version_id = ?) AS page_versions
            """,
            (accepted["document_id"], accepted["version_id"]),
        ).fetchone()
    assert dict(counts) == {"pages": 2, "page_versions": 2}


def test_failed_hidden_page_enablement_is_retryable_without_partial_page_facts(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "hidden-page-failure-source"},
        files={"file": ("hidden.pptx", build_minimal_presentation(), PPTX_MEDIA_TYPE)},
    ).json()
    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent((), (), (), ()),
    )
    monkeypatch.setattr(
        "pptextract.ingest_workflow.render_standard_pages",
        lambda *_args, pages, **_kwargs: (
            StandardPageRender(
                page_number=pages[0].page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=b"render",
            ),
        ),
    )
    assert run_once(settings) is True
    path = (
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}/source-pages/2/enable"
    )
    started = client.post(path, headers={"Idempotency-Key": "enable-that-fails"}).json()

    def unavailable_conversion(*_args: object, **_kwargs: object) -> object:
        raise OSError("converter remains unavailable")

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", unavailable_conversion)
    assert [run_once(settings) for _ in range(3)] == [True, True, True]

    failed = client.get(f"/api/v1/jobs/{started['job_id']}").json()
    assert failed["status"] == "failed"
    all_pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"]
    hidden = all_pages[1]
    assert hidden["enabled"] is False
    assert hidden["enablement"]["status"] == "failed"
    assert hidden["enablement"]["error"]["phase"] == "conversion"
    with connect(settings) as connection:
        partial = connection.execute(
            """
            SELECT enabled, source_content_json, fingerprint_sha256, render_sha256,
                   (SELECT COUNT(*) FROM page_versions
                    WHERE version_id = ? AND page_number = 2) AS page_versions
            FROM ingestion_page_results
            WHERE version_id = ? AND page_number = 2
            """,
            (accepted["version_id"], accepted["version_id"]),
        ).fetchone()
    assert dict(partial) == {
        "enabled": 0,
        "source_content_json": None,
        "fingerprint_sha256": None,
        "render_sha256": None,
        "page_versions": 0,
    }

    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent(("重试成功",), (), (), ()),
    )
    retried = client.post(path, headers={"Idempotency-Key": "retry-hidden-page"})
    assert retried.status_code == 202
    assert retried.json()["job_id"] != started["job_id"]
    assert run_once(settings) is True
    pending_pages = client.get("/api/v1/curation/pages").json()["pages"]
    assert [page["page_number"] for page in pending_pages] == [1, 2]


def test_expired_hidden_page_lease_clears_partial_page_facts(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    accepted = client.post(
        "/api/v1/documents",
        headers={"Idempotency-Key": "hidden-page-expired-source"},
        files={"file": ("hidden.pptx", build_minimal_presentation(), PPTX_MEDIA_TYPE)},
    ).json()
    monkeypatch.setattr(
        "pptextract.ingest_workflow.convert_page",
        lambda *_args, **_kwargs: NormalizedPageContent((), (), (), ()),
    )
    monkeypatch.setattr(
        "pptextract.ingest_workflow.render_standard_pages",
        lambda *_args, pages, **_kwargs: (
            StandardPageRender(
                page_number=pages[0].page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=b"render",
            ),
        ),
    )
    assert run_once(settings) is True
    path = (
        f"/api/v1/documents/{accepted['document_id']}"
        f"/versions/{accepted['version_id']}/source-pages/2/enable"
    )
    started = client.post(path, headers={"Idempotency-Key": "expired-hidden-page"}).json()

    for expected_attempt in (1, 2, 3):
        claim = claim_next_job(settings)
        assert claim is not None
        assert claim.attempts == expected_attempt
        with transaction(settings) as connection:
            connection.execute(
                """
                UPDATE jobs SET lease_expires_at = '2000-01-01T00:00:00+00:00'
                WHERE job_id = ?
                """,
                (started["job_id"],),
            )
            connection.execute(
                """
                UPDATE ingestion_page_results
                SET source_content_json = '{}', fingerprint_version = 1,
                    fingerprint_sha256 = 'partial'
                WHERE version_id = ? AND page_number = 2
                """,
                (accepted["version_id"],),
            )

    assert claim_next_job(settings) is None
    with connect(settings) as connection:
        page_result = connection.execute(
            """
            SELECT enabled, source_content_json, fingerprint_version, fingerprint_sha256
            FROM ingestion_page_results
            WHERE version_id = ? AND page_number = 2
            """,
            (accepted["version_id"],),
        ).fetchone()
    assert dict(page_result) == {
        "enabled": 0,
        "source_content_json": None,
        "fingerprint_version": None,
        "fingerprint_sha256": None,
    }
