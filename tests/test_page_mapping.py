from __future__ import annotations

from collections.abc import Iterator
from dataclasses import replace
from io import BytesIO
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from pptextract.api import create_app
from pptextract.config import Settings
from pptextract.conversion import NormalizedPageContent
from pptextract.db import transaction
from pptextract.object_store import LocalObjectStore
from pptextract.pptx_projection import SourcePage, list_source_pages
from pptextract.rendering import StandardPageRender
from pptextract.worker import run_once

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture
def system(tmp_path: Path) -> Iterator[tuple[TestClient, Settings]]:
    settings = replace(Settings.for_test(tmp_path), job_retry_base_seconds=0)
    with TestClient(create_app(settings)) as client:
        yield client, settings


def _presentation(*titles: str) -> bytes:
    presentation = Presentation()
    for title in titles:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _replace_title(source: bytes, old_title: str, new_title: str) -> bytes:
    presentation = Presentation(BytesIO(source))
    slide = next(
        slide for slide in presentation.slides if slide.shapes.title.text == old_title
    )
    slide.shapes.title.text = new_title
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _install_toolchain(monkeypatch: pytest.MonkeyPatch) -> None:
    def convert(source: bytes, page: SourcePage) -> NormalizedPageContent:
        presentation = Presentation(BytesIO(source))
        title = presentation.slides[page.page_number - 1].shapes.title.text
        return NormalizedPageContent((title,), (), (), ())

    def render(
        _source: bytes, *, toolchain: object, pages: tuple[SourcePage, ...]
    ) -> tuple[StandardPageRender, ...]:
        del toolchain
        page_number = pages[0].page_number
        return (
            StandardPageRender(
                page_number=page_number,
                media_type="image/png",
                dpi=144,
                width_px=10,
                height_px=10,
                data=f"render-{page_number}".encode(),
            ),
        )

    monkeypatch.setattr("pptextract.ingest_workflow.convert_page", convert)
    monkeypatch.setattr("pptextract.ingest_workflow.render_standard_pages", render)


def _upload(
    client: TestClient,
    source: bytes,
    *,
    key: str,
    document_id: str | None = None,
) -> dict[str, str]:
    route = (
        "/api/v1/documents"
        if document_id is None
        else f"/api/v1/documents/{document_id}/versions"
    )
    response = client.post(
        route,
        headers={"Idempotency-Key": key},
        files={"file": (f"{key}.pptx", source, PPTX_MEDIA_TYPE)},
    )
    assert response.status_code == 202
    return response.json()


def _mapping_route(document_id: str, version_id: str) -> str:
    return f"/api/v1/documents/{document_id}/versions/{version_id}/page-mapping"


def test_duplicate_fingerprints_require_editable_atomic_mapping_decisions(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    baseline_source = _presentation("重复内容", "稳定锚点", "将变成重复内容")
    baseline = _upload(client, baseline_source, key="mapping-baseline")
    assert run_once(settings) is True
    baseline_pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"]
    baseline_by_number = {page["page_number"]: page for page in baseline_pages}

    incoming = _upload(
        client,
        _replace_title(baseline_source, "将变成重复内容", "重复内容"),
        key="mapping-duplicate",
        document_id=baseline["document_id"],
    )
    assert run_once(settings) is True

    version = client.get(
        f"/api/v1/documents/{baseline['document_id']}"
        f"/versions/{incoming['version_id']}"
    ).json()
    task = client.get(f"/api/v1/jobs/{incoming['job_id']}").json()
    document = client.get(f"/api/v1/documents/{baseline['document_id']}").json()
    assert version["status"] == "awaiting_mapping"
    assert task["status"] == "requires_action"
    assert document["current_version_id"] == baseline["version_id"]
    assert {
        page["version_id"]
        for page in client.get(
            "/api/v1/curation/pages", params={"review_status": "all"}
        ).json()["pages"]
    } == {baseline["version_id"]}
    processing = next(
        runway
        for runway in client.get("/api/v1/app/bootstrap").json()["runways"]
        if runway["id"] == "processing"
    )
    assert processing["documents"] == [
        {
            "document_id": baseline["document_id"],
            "version_id": incoming["version_id"],
            "title": "mapping-duplicate.pptx",
            "status": "requires_action",
            "status_label": "需要页对应",
                "action": {
                "label": "处理页对应",
                "href": (
                    f"/documents/{baseline['document_id']}/versions/"
                    f"{incoming['version_id']}/page-mapping"
                    ),
                },
                "rendering_warnings": {
                    "total": 3,
                    "pages": 3,
                    "unconfirmed": 3,
                    "unconfirmed_pages": 3,
                },
            }
    ]

    route = _mapping_route(baseline["document_id"], incoming["version_id"])
    response = client.get(route)
    assert response.status_code == 200
    etag = response.headers["ETag"]
    workspace = response.json()
    assert workspace["remaining_cases"] == 2
    assert workspace["current_version"] == {
        "version_id": baseline["version_id"],
        "still_serving": True,
    }
    assert [case["source_page"]["page_number"] for case in workspace["cases"]] == [1, 3]
    assert all(case["kind"] == "duplicate_fingerprint" for case in workspace["cases"])
    assert all(case["decision"] is None for case in workspace["cases"])
    assert all(case["candidates"] for case in workspace["cases"])
    assert all(
        "adjacent_confirmed" in candidate and "relative_order" in candidate
        for case in workspace["cases"]
        for candidate in case["candidates"]
    )
    assert workspace["cases"][0]["source_page"]["slide_id"] is not None
    assert workspace["cases"][0]["source_page"]["standard_render"]["url"].endswith(
        "/source-pages/1/render"
    )
    source_render = client.get(
        workspace["cases"][0]["source_page"]["standard_render"]["url"]
    )
    candidate_render = client.get(
        workspace["cases"][0]["candidates"][0]["standard_render"]["url"]
    )
    assert source_render.status_code == candidate_render.status_code == 200
    assert source_render.headers["content-type"] == "image/png"

    busy = client.post(
        f"/api/v1/documents/{baseline['document_id']}/versions",
        headers={"Idempotency-Key": "mapping-still-serial"},
        files={
            "file": (
                "mapping-still-serial.pptx",
                _presentation("不能越过人工门禁"),
                PPTX_MEDIA_TYPE,
            )
        },
    )
    assert busy.status_code == 409
    assert busy.json()["error"]["code"] == "document_busy"

    unresolved = client.post(
        f"{route}/confirm",
        headers={"If-Match": etag, "X-Actor-ID": "operator-li"},
    )
    assert unresolved.status_code == 409
    assert unresolved.json()["error"]["code"] == "mapping_incomplete"

    first_case, second_case = workspace["cases"]
    reused_page_id = baseline_by_number[1]["page_id"]
    first_save = client.put(
        f"{route}/cases/{first_case['case_id']}",
        headers={"If-Match": etag, "X-Actor-ID": "operator-li"},
        json={"decision": "reuse", "page_id": reused_page_id},
    )
    assert first_save.status_code == 200
    etag_after_first = first_save.headers["ETag"]

    stale = client.put(
        f"{route}/cases/{second_case['case_id']}",
        headers={"If-Match": etag, "X-Actor-ID": "operator-wang"},
        json={"decision": "new"},
    )
    assert stale.status_code == 412
    assert stale.json()["error"]["code"] == "mapping_precondition_failed"
    assert stale.headers["ETag"] == etag_after_first

    occupied = client.put(
        f"{route}/cases/{second_case['case_id']}",
        headers={"If-Match": etag_after_first, "X-Actor-ID": "operator-li"},
        json={"decision": "reuse", "page_id": reused_page_id},
    )
    assert occupied.status_code == 409
    assert occupied.json()["error"]["code"] == "mapping_candidate_occupied"

    second_save = client.put(
        f"{route}/cases/{second_case['case_id']}",
        headers={"If-Match": etag_after_first, "X-Actor-ID": "operator-li"},
        json={"decision": "new"},
    )
    assert second_save.status_code == 200
    etag_after_second = second_save.headers["ETag"]
    assert second_save.json()["remaining_cases"] == 0

    first_changed = client.put(
        f"{route}/cases/{first_case['case_id']}",
        headers={"If-Match": etag_after_second, "X-Actor-ID": "operator-li"},
        json={"decision": "new"},
    )
    assert first_changed.status_code == 200
    etag_after_change = first_changed.headers["ETag"]
    second_changed = client.put(
        f"{route}/cases/{second_case['case_id']}",
        headers={"If-Match": etag_after_change, "X-Actor-ID": "operator-li"},
        json={"decision": "reuse", "page_id": reused_page_id},
    )
    assert second_changed.status_code == 200
    final_etag = second_changed.headers["ETag"]

    with transaction(settings) as connection:
        missing_render = connection.execute(
            """
            SELECT render_sha256 FROM ingestion_page_results
            WHERE version_id = ? AND page_number = 1
            """,
            (incoming["version_id"],),
        ).fetchone()["render_sha256"]
    store = LocalObjectStore(settings.object_store_path)
    store.path_for(missing_render).unlink()
    unavailable = client.get(route).json()
    assert unavailable["impact_summary"]["evidence_errors"] >= 1
    assert unavailable["can_confirm"] is False
    blocked_confirmation = client.post(
        f"{route}/confirm",
        headers={"If-Match": final_etag, "X-Actor-ID": "operator-li"},
    )
    assert blocked_confirmation.status_code == 409
    assert blocked_confirmation.json()["error"]["code"] == "mapping_evidence_unavailable"
    store.put(b"render-1")

    confirmed = client.post(
        f"{route}/confirm",
        headers={"If-Match": final_etag, "X-Actor-ID": "operator-li"},
    )
    assert confirmed.status_code == 200
    assert confirmed.json()["status"] == "ready"
    assert confirmed.json()["summary"] == {
        "reused_unchanged": 2,
        "reused_changed": 0,
        "created_new": 1,
        "soft_deleted": 1,
    }

    current = client.get(f"/api/v1/documents/{baseline['document_id']}").json()
    assert current["current_version_id"] == incoming["version_id"]
    pages = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"]
    duplicate_pages = [page for page in pages if page["title"] == "重复内容"]
    assert len(duplicate_pages) == 2
    assert len({page["page_id"] for page in duplicate_pages}) == 2
    assert len({page["chunk_id"] for page in duplicate_pages}) == 2

    frozen = client.put(
        f"{route}/cases/{first_case['case_id']}",
        headers={"If-Match": confirmed.headers["ETag"], "X-Actor-ID": "operator-li"},
        json={"decision": "reuse", "page_id": reused_page_id},
    )
    assert frozen.status_code == 409
    assert frozen.json()["error"]["code"] == "mapping_frozen"


def test_slide_id_conflict_exposes_each_candidate_instead_of_silently_mapping(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    baseline_source = _presentation("历史甲", "稳定锚点", "历史乙")
    baseline = _upload(client, baseline_source, key="slide-conflict-baseline")
    assert run_once(settings) is True
    with transaction(settings) as connection:
        connection.execute(
            """
            UPDATE ingestion_page_results SET source_slide_id = 900
            WHERE version_id = ? AND page_number IN (1, 3)
            """,
            (baseline["version_id"],),
        )

    incoming_source = _replace_title(
        _replace_title(baseline_source, "历史甲", "变化甲"), "历史乙", "变化乙"
    )
    incoming = _upload(
        client,
        incoming_source,
        key="slide-conflict-incoming",
        document_id=baseline["document_id"],
    )

    def conflicting_manifest(source: bytes) -> tuple[SourcePage, ...]:
        return tuple(
            replace(page, source_slide_id=900)
            if page.page_number in {1, 3}
            else page
            for page in list_source_pages(source)
        )

    monkeypatch.setattr(
        "pptextract.ingest_workflow.list_source_pages", conflicting_manifest
    )
    assert run_once(settings) is True

    route = _mapping_route(baseline["document_id"], incoming["version_id"])
    workspace = client.get(route).json()
    assert [case["kind"] for case in workspace["cases"]] == [
        "slide_id_conflict",
        "slide_id_conflict",
    ]
    assert [len(case["candidates"]) for case in workspace["cases"]] == [2, 2]
    assert {
        candidate["page_id"]
        for candidate in workspace["cases"][0]["candidates"]
    } == {
        candidate["page_id"]
        for candidate in workspace["cases"][1]["candidates"]
    }
    assert client.get(f"/api/v1/jobs/{incoming['job_id']}").json()["status"] == (
        "requires_action"
    )


def test_one_historical_identity_with_two_known_contents_requires_human_choice(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    first_source = _presentation("历史内容甲")
    first = _upload(client, first_source, key="multi-history-first")
    assert run_once(settings) is True
    original_page_id = client.get(
        "/api/v1/curation/pages", params={"review_status": "all"}
    ).json()["pages"][0]["page_id"]

    second = _upload(
        client,
        _replace_title(first_source, "历史内容甲", "历史内容乙"),
        key="multi-history-second",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True
    assert client.get(
        f"/api/v1/documents/{first['document_id']}/versions/{second['version_id']}"
    ).json()["status"] == "ready"

    third = _upload(
        client,
        _presentation("历史内容甲", "历史内容乙"),
        key="multi-history-third",
        document_id=first["document_id"],
    )
    assert run_once(settings) is True

    route = _mapping_route(first["document_id"], third["version_id"])
    workspace = client.get(route).json()
    assert workspace["remaining_cases"] == 2
    assert [case["kind"] for case in workspace["cases"]] == [
        "multiple_candidates",
        "multiple_candidates",
    ]
    assert all(
        [candidate["page_id"] for candidate in case["candidates"]]
        == [original_page_id]
        for case in workspace["cases"]
    )
    assert client.get(f"/api/v1/jobs/{third['job_id']}").json()["status"] == (
        "requires_action"
    )


def test_new_duplicate_content_is_not_silently_mapped_by_unique_slide_ids(
    system: tuple[TestClient, Settings], monkeypatch: pytest.MonkeyPatch
) -> None:
    client, settings = system
    _install_toolchain(monkeypatch)
    baseline_source = _presentation("原内容甲", "原内容乙")
    baseline = _upload(client, baseline_source, key="new-duplicate-baseline")
    assert run_once(settings) is True

    duplicate_source = _replace_title(
        _replace_title(baseline_source, "原内容甲", "共同新内容"),
        "原内容乙",
        "共同新内容",
    )
    incoming = _upload(
        client,
        duplicate_source,
        key="new-duplicate-incoming",
        document_id=baseline["document_id"],
    )
    assert run_once(settings) is True

    route = _mapping_route(baseline["document_id"], incoming["version_id"])
    workspace = client.get(route).json()
    assert [case["kind"] for case in workspace["cases"]] == [
        "duplicate_fingerprint",
        "duplicate_fingerprint",
    ]
    assert [len(case["candidates"]) for case in workspace["cases"]] == [1, 1]
    assert client.get(f"/api/v1/jobs/{incoming['job_id']}").json()["status"] == (
        "requires_action"
    )
